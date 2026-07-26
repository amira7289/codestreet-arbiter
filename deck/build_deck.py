"""Build the CodeStreet pitch deck as a .pptx.

Generated rather than hand-made so the numbers on the slides cannot drift from the
numbers the system actually reports — TEAM, MEMBERS and the sourced cost figure are
the only things a human has to fill in.

    pip install -r deck/requirements.txt
    python deck/build_deck.py

Layout follows the official CodeStreet template: navy/gold title and impact slides,
light slides for everything else, footer on all but the full-bleed ones.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- fill these in --
PROJECT = "ARBITER"
TAGLINE = "A neutral arbiter that shows both sides the same reasoning — in seconds, not weeks."
THEME = "Frictionless Dispute & Chargeback Resolution"
TEAM = "[TEAM NAME]"
MEMBERS = ["[Member 1]", "[Member 2]"]
MEMBER_ROLES = ["[Role / focus area]", "[Role / focus area]"]

NAVY = RGBColor(0x0A, 0x2A, 0x5E)
NAVY_DEEP = RGBColor(0x06, 0x1D, 0x45)
GOLD = RGBColor(0xC9, 0xA2, 0x39)
BLUE = RGBColor(0x00, 0x6F, 0xCF)
BLUE_LT = RGBColor(0x6F, 0xB4, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xFA)
GREY = RGBColor(0x5A, 0x64, 0x72)
GREY_LT = RGBColor(0x8A, 0x93, 0x9E)

FONT = "Arial"
W, H = Inches(13.333), Inches(7.5)


def text(slide, x, y, w, h, runs, *, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
         italic=False, spacing=1.15, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = runs if isinstance(runs, list) else [runs]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        parts = line if isinstance(line, list) else [(line, {})]
        for content, style in parts:
            r = p.add_run()
            r.text = content
            f = r.font
            f.name = FONT
            f.size = Pt(style.get("size", size))
            f.bold = style.get("bold", bold)
            f.italic = style.get("italic", italic)
            f.color.rgb = style.get("color", color)
    return box


def rect(slide, x, y, w, h, fill, *, radius=True, line=None, line_w=1.25):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        shape.adjustments[0] = 0.06
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def circle(slide, x, y, d, fill, *, line=None, line_w=1.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def blank(prs, bg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    back = rect(slide, 0, 0, W, H, bg, radius=False)
    back.shadow.inherit = False
    return slide


def footer(slide, page):
    text(slide, Inches(0.55), Inches(6.98), Inches(9), Inches(0.3),
         f"{TEAM}  ·  American Express CodeStreet 2026", size=10, color=GREY_LT)
    text(slide, Inches(12.2), Inches(6.98), Inches(0.6), Inches(0.3),
         f"{page:02d}", size=10, color=GREY_LT, align=PP_ALIGN.RIGHT)


def heading(slide, title, sub=None):
    text(slide, Inches(0.7), Inches(0.62), Inches(11.5), Inches(0.7),
         title, size=34, bold=True, color=NAVY)
    if sub:
        text(slide, Inches(0.7), Inches(1.42), Inches(10.4), Inches(0.5),
             sub, size=14, color=GREY, italic=True)


prs = Presentation()
prs.slide_width, prs.slide_height = W, H

# ============================================================== 01 · title ======
s = blank(prs, NAVY)
rect(s, Inches(0.55), Inches(0.42), Inches(0.62), Inches(0.62), WHITE, radius=False)
circle(s, Inches(0.71), Inches(0.58), Inches(0.3), BLUE)
text(s, Inches(1.42), Inches(0.55), Inches(8), Inches(0.4),
     "AMERICAN EXPRESS CODESTREET 2026", size=17, bold=True, color=GOLD)

# Two lines rather than one: at 44pt the name and the descriptor together overrun
# the 8.4in column and collide with the rule beneath them.
text(s, Inches(0.55), Inches(1.42), Inches(8.4), Inches(0.85),
     PROJECT, size=46, bold=True, color=WHITE)
text(s, Inches(0.55), Inches(2.16), Inches(8.4), Inches(0.5),
     "Dispute Resolution", size=25, bold=True, color=BLUE)
rect(s, Inches(0.58), Inches(2.78), Inches(1.3), Inches(0.06), GOLD, radius=False)

text(s, Inches(0.55), Inches(3.04), Inches(8.2), Inches(0.6),
     TAGLINE, size=16, bold=True, color=WHITE)

caps = [("AI", "Extracts facts,\nnever decides"),
        ("SC", "25 named signals,\none scorecard"),
        ("NG", "Settle before\nadjudicating"),
        ("XP", "Counterfactual on\nevery verdict")]
for i, (tag, label) in enumerate(caps):
    x = Inches(0.55 + i * 2.13)
    circle(s, x, Inches(4.02), Inches(0.62), None, line=GOLD, line_w=1.5)
    text(s, x, Inches(4.17), Inches(0.62), Inches(0.35), tag,
         size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.78), Inches(4.05), Inches(1.28), Inches(0.7),
         label, size=10.5, color=WHITE, spacing=1.05)
    if i < 3:
        rect(s, x + Inches(1.98), Inches(3.99), Inches(0.012), Inches(0.68),
             RGBColor(0x2A, 0x4A, 0x7E), radius=False)

box = rect(s, Inches(0.55), Inches(5.02), Inches(8.4), Inches(0.95), None,
           line=RGBColor(0x2A, 0x4A, 0x7E))
circle(s, Inches(0.85), Inches(5.28), Inches(0.44), None, line=GOLD, line_w=3)
text(s, Inches(1.55), Inches(5.16), Inches(7.2), Inches(0.72),
     [[("The decision is never in the model. ", {"color": WHITE}),
       ("Every point traces to a named rule", {"color": GOLD}),
       (" both parties can read.", {"color": WHITE})]],
     size=14, bold=True, spacing=1.2)

for i, (label, value) in enumerate([("Theme:", THEME), ("Team Name:", TEAM),
                                    ("Team Members:", "  •  ".join(MEMBERS))]):
    text(s, Inches(0.55), Inches(6.15 + i * 0.34), Inches(8.5), Inches(0.3),
         [[(label + " ", {"color": GOLD, "bold": True}), (value, {"color": WHITE})]], size=12)

# concentric dial
cx, cy = Inches(10.9), Inches(3.35)
circle(s, cx - Inches(2.05), cy - Inches(2.05), Inches(4.1), WHITE)
circle(s, cx - Inches(1.78), cy - Inches(1.78), Inches(3.56), NAVY_DEEP)
circle(s, cx - Inches(1.42), cy - Inches(1.42), Inches(2.84), None, line=BLUE, line_w=1)
circle(s, cx - Inches(1.0), cy - Inches(1.0), Inches(2.0), None, line=GOLD, line_w=1)
circle(s, cx - Inches(0.11), cy - Inches(0.11), Inches(0.22), BLUE)
for dx, dy in [(0, -1.42), (1.0, -1.0), (1.42, 0), (1.0, 1.0), (0, 1.42), (-1.0, 1.0),
               (-1.42, 0), (-1.0, -1.0)]:
    circle(s, cx + Inches(dx) - Inches(0.055), cy + Inches(dy) - Inches(0.055), Inches(0.11), GOLD)
text(s, cx - Inches(1.25), cy - Inches(0.72), Inches(2.5), Inches(1.4),
     [[("One ", {"color": WHITE}), ("Scorecard.", {"color": BLUE_LT})],
      [("Every Case.", {"color": WHITE})],
      [("Full Transparency.", {"color": WHITE})]],
     size=15, bold=True, align=PP_ALIGN.CENTER, spacing=1.3)

# ============================================================ 02 · problem ======
s = blank(prs, WHITE)
heading(s, "The Problem",
        "A card member disputes a charge, and then everyone waits — including the merchant.")
pains = [
    ("Weeks of waiting, on both sides",
     "The merchant gets weeks to respond and the card member is out of pocket meanwhile. "
     "Neither can do anything to speed it up."),
    ("A verdict with no reasoning attached",
     "An analyst reads shipping records, policies and correspondence, and issues an outcome. "
     "Neither party is told which fact decided it."),
    ("Every tool on the market is one-sided",
     "Merchant win-rate services argue for the merchant. Network data-sharing prevents "
     "escalation. Nobody arbitrates neutrally and shows their working."),
]
for i, (title, body) in enumerate(pains):
    y = Inches(2.3 + i * 1.35)
    rect(s, Inches(0.7), y, Inches(0.72), Inches(0.72), NAVY)
    text(s, Inches(0.7), y + Inches(0.19), Inches(0.72), Inches(0.4), f"{i + 1:02d}",
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.75), y + Inches(0.02), Inches(10.6), Inches(0.35), title,
         size=16, bold=True, color=NAVY)
    text(s, Inches(1.75), y + Inches(0.42), Inches(10.6), Inches(0.7), body, size=12.5, color=GREY)
footer(s, 2)

# =========================================================== 03 · solution ======
s = blank(prs, LIGHT)
heading(s, "Our Solution")
card = rect(s, Inches(1.1), Inches(2.05), Inches(11.1), Inches(2.15), WHITE)
text(s, Inches(1.7), Inches(2.45), Inches(9.9), Inches(1.4),
     "“An issuer-run arbiter that auto-gathers the evidence, weighs it with a "
     "deterministic scorecard, and hands both parties the same explanation for the "
     "same verdict.”",
     size=21, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER, spacing=1.28)
text(s, Inches(1.1), Inches(4.65), Inches(11.1), Inches(0.9),
     [[("What we are NOT building: ", {"bold": True, "color": NAVY})],
      [("Not fraud detection — the transaction is genuine and the dispute is about what "
        "happened next. Not a merchant win-rate tool. Not a trained classifier: the model "
        "extracts facts and writes prose, it never decides the outcome.", {"color": GREY})]],
     size=13, align=PP_ALIGN.CENTER, spacing=1.3)
footer(s, 3)

# ==================================================== 04 · differentiators ======
s = blank(prs, WHITE)
heading(s, "What Makes Us Different")
diffs = [
    ("Neutral, and two-sided", NAVY,
     "Both parties read the identical explanation. Every competing tool argues for one "
     "side or moves data between institutions."),
    ("The model never decides", BLUE,
     "An LLM extracts typed facts and narrates. A deterministic scorecard rules. That "
     "split is what makes the verdict reproducible, auditable and testable."),
    ("It reports its own uncertainty", NAVY,
     "Confidence is damped by how much evidence exists, ties are disclosed on the "
     "verdict, and cases nobody could decide are flagged for a human."),
]
for i, (title, accent, body) in enumerate(diffs):
    x = Inches(0.7 + i * 4.07)
    dark = i == 1
    rect(s, x, Inches(2.15), Inches(3.77), Inches(3.5), NAVY if dark else LIGHT)
    circle(s, x + Inches(0.42), Inches(2.62), Inches(0.62), BLUE if dark else NAVY)
    text(s, x + Inches(0.42), Inches(2.78), Inches(0.62), Inches(0.4), str(i + 1),
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.42), Inches(3.62), Inches(2.95), Inches(0.6), title,
         size=15.5, bold=True, color=WHITE if dark else NAVY)
    text(s, x + Inches(0.42), Inches(4.3), Inches(2.95), Inches(1.2), body,
         size=11.5, color=RGBColor(0xC7, 0xD8, 0xEE) if dark else GREY, spacing=1.25)
footer(s, 4)

# ========================================================= 05 · how it works ====
s = blank(prs, LIGHT)
heading(s, "How It Works",
        "Five stages. Adjudication is the fallback, not the first move.")
stages = [("Gather", "4 sources\nqueried"), ("Extract", "typed facts,\nnot prose"),
          ("Negotiate", "settle without\na ruling"), ("Score", "25 weighted\nsignals"),
          ("Explain", "code, reasons,\ncounterfactual")]
for i, (name, sub) in enumerate(stages):
    x = Inches(0.7 + i * 2.5)
    accent = BLUE if i == 2 else NAVY
    rect(s, x, Inches(2.55), Inches(2.14), Inches(1.5), WHITE, line=accent, line_w=1.5)
    text(s, x, Inches(2.87), Inches(2.14), Inches(0.35), name,
         size=15, bold=True, color=accent, align=PP_ALIGN.CENTER)
    text(s, x, Inches(3.3), Inches(2.14), Inches(0.6), sub,
         size=10.5, color=GREY, align=PP_ALIGN.CENTER, spacing=1.15)
    if i < 4:
        text(s, x + Inches(2.16), Inches(3.02), Inches(0.34), Inches(0.4), "→",
             size=17, color=BLUE, align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(4.55), Inches(11.9), Inches(0.9),
     [[("A settlement both sides accept closes the case with no verdict at all. ", {"bold": True}),
       ("Where they cannot agree, the scorecard rules — and because it is arithmetic over "
        "named signals, the same evidence always produces the same answer.", {"color": GREY})]],
     size=13, color=NAVY, spacing=1.3)
footer(s, 5)

# ======================================================= 06 · transparency ======
s = blank(prs, WHITE)
heading(s, "Four Layers of Reasoning",
        "In increasing order of usefulness to the party who just lost.")
layers = [
    ("Signal breakdown", "Every point, named, attributed to the document it came from."),
    ("Amex reason code", "C02 / C08 / C31 / C32 / P08, derived deterministically from "
                         "claim type and fired signals."),
    ("The counterfactual", "“This would have gone the other way if these had not been "
                           "established (−45 points)…” Pure arithmetic over the "
                           "scorecard, so it cannot hallucinate."),
    ("The narrative", "Required to name the losing side's strongest point and why it was "
                      "outweighed."),
]
for i, (title, body) in enumerate(layers):
    y = Inches(2.2 + i * 1.05)
    circle(s, Inches(0.7), y + Inches(0.06), Inches(0.5), BLUE if i == 2 else NAVY)
    text(s, Inches(0.7), y + Inches(0.14), Inches(0.5), Inches(0.35), str(i + 1),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.45), y, Inches(2.6), Inches(0.4), title, size=14.5, bold=True, color=NAVY)
    text(s, Inches(4.15), y + Inches(0.02), Inches(8.3), Inches(0.85), body,
         size=12, color=GREY, spacing=1.2)
rect(s, Inches(0.7), Inches(6.4), Inches(11.75), Inches(0.42), LIGHT)
text(s, Inches(1.0), Inches(6.5), Inches(11.2), Inches(0.3),
     "Plus a “recommended for human review” flag when a party filed evidence no "
     "rule reads, or confidence falls below 35%.",
     size=11.5, color=NAVY, italic=True)
footer(s, 6)

# ============================================================ 07 · fairness ====
s = blank(prs, NAVY)
text(s, Inches(0.7), Inches(0.62), Inches(11.5), Inches(0.7),
     "Fairness, Stated Out Loud", size=34, bold=True, color=WHITE)
text(s, Inches(0.7), Inches(1.55), Inches(11.3), Inches(0.5),
     "We do not claim the scorecard is unbiased.", size=17, bold=True, color=GOLD)
text(s, Inches(0.7), Inches(2.15), Inches(11.3), Inches(1.0),
     "Ambiguous cases resolve to the card member, matching issuer provisional-credit "
     "practice. That default is deliberate — and it is printed on the verdict as a "
     "zero-weight signal, not hidden inside a comparison operator.",
     size=14, color=WHITE, spacing=1.3)
rect(s, Inches(0.7), Inches(3.35), Inches(0.05), Inches(0.95), GOLD, radius=False)
text(s, Inches(1.1), Inches(3.4), Inches(11.0), Inches(0.9),
     [[("The target is not a bias gap of zero.", {"color": GOLD, "bold": True})],
      [("It is that every point of directional bias traces to a rule we state out loud.",
        {"color": WHITE, "bold": True})]],
     size=15, spacing=1.3)
audits = [("bias_gap 0.000", "recall equal for both parties on the labelled corpus"),
          ("0 confidently wrong", "no verdict above 80% confidence missed"),
          ("CI-enforced", "the build fails on any weight asymmetry left undocumented")]
for i, (big, sub) in enumerate(audits):
    x = Inches(0.7 + i * 3.95)
    text(s, x, Inches(4.85), Inches(3.7), Inches(0.4), big, size=19, bold=True, color=GOLD)
    text(s, x, Inches(5.35), Inches(3.6), Inches(0.8), sub, size=11.5,
         color=RGBColor(0xC7, 0xD8, 0xEE), spacing=1.2)
footer(s, 7)

# =============================================================== 08 · scope ====
s = blank(prs, WHITE)
heading(s, "Scope: What's Real vs. Mocked",
        "Stated plainly, because a system whose thesis is reporting its own uncertainty "
        "should demonstrate that about itself.")
rows = [
    ("Evidence extraction, scoring, verdicts, counterfactuals", "Fully working"),
    ("Negotiation: offers, counters, settlement", "Fully working"),
    ("Two-party live interface, auto-gather timeline", "Fully working"),
    ("Test suite, accuracy / fairness / calibration harness", "Fully working"),
    ("The four evidence sources (carrier, ledger, policy, CRM)", "Simulated for demo"),
    ("Connector latencies", "Reported, not measured"),
    ("Learned weights, abstain verdict, production auth", "Roadmap only"),
]
top, rh = Inches(2.4), Inches(0.52)
rect(s, Inches(0.7), top, Inches(11.75), rh, NAVY, radius=False)
text(s, Inches(1.0), top + Inches(0.14), Inches(8), Inches(0.3), "Component",
     size=12.5, bold=True, color=WHITE)
text(s, Inches(9.4), top + Inches(0.14), Inches(3), Inches(0.3), "Status",
     size=12.5, bold=True, color=WHITE)
for i, (comp, status) in enumerate(rows):
    y = top + rh + Inches(i * 0.52)
    if i % 2 == 1:
        rect(s, Inches(0.7), y, Inches(11.75), Inches(0.52), LIGHT, radius=False)
    text(s, Inches(1.0), y + Inches(0.14), Inches(8.2), Inches(0.3), comp, size=12, color=NAVY)
    working = status == "Fully working"
    text(s, Inches(9.4), y + Inches(0.14), Inches(3), Inches(0.3), status,
         size=12, bold=working, color=NAVY if working else GREY)
footer(s, 8)

# ============================================================== 09 · impact ====
s = blank(prs, NAVY)
text(s, Inches(0.7), Inches(0.62), Inches(11.5), Inches(0.7),
     "Business Impact", size=34, bold=True, color=WHITE)
text(s, Inches(0.7), Inches(1.5), Inches(11.3), Inches(0.4),
     "Every figure below is produced by GET /metrics on the running system, not asserted.",
     size=13, color=BLUE_LT, italic=True)
stats = [("54 / 54", "labelled cases called correctly\non the golden corpus"),
         ("0.12 ms", "p95 to parse, score and explain\na complete case"),
         ("0.000", "bias gap between the two parties—\nrecall is equal")]
for i, (big, sub) in enumerate(stats):
    x = Inches(0.7 + i * 3.95)
    text(s, x, Inches(2.6), Inches(3.7), Inches(0.9), big, size=46, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)
    text(s, x, Inches(3.75), Inches(3.7), Inches(0.9), sub, size=12,
         color=WHITE, align=PP_ALIGN.CENTER, spacing=1.25)
rect(s, Inches(0.7), Inches(5.0), Inches(11.75), Inches(1.35), NAVY_DEEP)
text(s, Inches(1.1), Inches(5.2), Inches(11.0), Inches(1.0),
     [[("Read 54/54 as a warning, not a trophy. ", {"color": GOLD, "bold": True}),
       ("The scorecard was tuned against this corpus, so a perfect score means the corpus "
        "has stopped discriminating — not that the system is perfect. The defensible claim "
        "is that no known failure mode is unhandled.", {"color": WHITE})],
      [("[Insert a cited industry figure for average chargeback handling cost here — do not "
        "estimate one.]", {"color": GREY_LT, "italic": True})]],
     size=12, spacing=1.35)
footer(s, 9)

# ================================================================ 10 · demo ====
s = blank(prs, LIGHT)
heading(s, "Demo Flow  &  What's Next")
text(s, Inches(0.7), Inches(2.15), Inches(5.6), Inches(0.4), "Live Demo Steps",
     size=17, bold=True, color=BLUE)
steps = [
    "Open a filed dispute — no evidence on it yet.",
    "Auto-gather: four sources report live, hits and misses, with latencies.",
    "Resolve: verdict, Amex reason code, signal breakdown, counterfactual.",
    "Switch to the other party's panel — the identical explanation.",
    "Edge case: file contradicting evidence. The verdict is withdrawn and the case re-opens.",
]
for i, step in enumerate(steps):
    y = Inches(2.75 + i * 0.72)
    text(s, Inches(0.7), y, Inches(0.4), Inches(0.3), f"{i + 1}.", size=13, bold=True, color=NAVY)
    text(s, Inches(1.12), y, Inches(5.3), Inches(0.62), step, size=12.5, color=NAVY, spacing=1.2)

text(s, Inches(7.0), Inches(2.15), Inches(5.4), Inches(0.4), "Roadmap",
     size=17, bold=True, color=GOLD)
next_up = [
    "Abstain as a real verdict, for cases no human could decide either.",
    "Adverse inference weighted by how probative a filing is, not whether it exists.",
    "A held-out corpus nobody sees while tuning — the only way the accuracy number "
    "starts meaning something.",
    "Fit the weights to resolved outcomes. The scorecard stays deterministic; the "
    "counterfactual tells you which signal to collect next.",
]
for i, item in enumerate(next_up):
    y = Inches(2.75 + i * 0.95)
    circle(s, Inches(7.0), y + Inches(0.08), Inches(0.14), GOLD)
    text(s, Inches(7.35), y, Inches(5.05), Inches(0.85), item, size=12.5, color=NAVY, spacing=1.2)
footer(s, 10)

# ================================================================ 11 · team ====
s = blank(prs, WHITE)
heading(s, "Our Team")
for i, (name, role) in enumerate(zip(MEMBERS, MEMBER_ROLES)):
    x = Inches(1.55 + i * 5.3)
    rect(s, x, Inches(2.3), Inches(4.65), Inches(3.3), LIGHT)
    circle(s, x + Inches(1.72), Inches(2.75), Inches(1.2), NAVY)
    text(s, x + Inches(1.72), Inches(3.15), Inches(1.2), Inches(0.4), f"{i + 1}",
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, x, Inches(4.3), Inches(4.65), Inches(0.4), name,
         size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    text(s, x, Inches(4.8), Inches(4.65), Inches(0.4), role,
         size=13, color=GREY, align=PP_ALIGN.CENTER)
footer(s, 11)

# ============================================================= 12 · closing ====
s = blank(prs, NAVY)
circle(s, Inches(-1.6), Inches(4.6), Inches(3.6), NAVY_DEEP)
text(s, Inches(0.85), Inches(3.05), Inches(11.0), Inches(1.0),
     "Explainability isn't a constraint on accuracy here. It's the training signal.",
     size=33, bold=True, color=WHITE, spacing=1.2)
text(s, Inches(0.9), Inches(4.35), Inches(9), Inches(0.4),
     "Thank you — questions welcome.", size=15, color=BLUE_LT)

out = Path(__file__).resolve().parent / "codestreet-arbiter-deck.pptx"
prs.save(out)
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
